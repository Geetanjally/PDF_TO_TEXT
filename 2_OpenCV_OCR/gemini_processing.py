import os
import google.generativeai as genai
from pptx import Presentation

# -----------------------------
# 1. Gemini CONFIG
# -----------------------------
genai.configure(api_key="AIzaSyAzHf66I6a1uHUbC1-PnFCK6KyBUZTOJYI")
MODEL_NAME = "models/gemini-2.5-flash"
model = genai.GenerativeModel(MODEL_NAME)


def clean_with_gemini(raw_text):
    prompt = f"""
You are a professional text cleaner.

Clean the following messy OCR extracted text:
- Fix spelling mistakes
- Fix sentence breaks
- Remove unwanted symbols or OCR noise
- Preserve bullet points
- Preserve headings
- Convert into clean, structured text

OCR Text:
{raw_text}

Return ONLY the cleaned text.
"""

    try:
        # Use global model
        response = model.generate_content(
            prompt,
            generation_config={
                "temperature": 0.2,
                "top_p": 1,
                "top_k": 40,
                "max_output_tokens": 8000,
            }
        )

        return response.text.strip()

    except Exception as e:
        print("❌ Gemini Error:", e)
        return f"[ERROR] {e}"

# -----------------------------
# 3. CREATE PPT FROM CLEAN TEXT
# -----------------------------
from pptx import Presentation
from pptx.util import Pt, Inches

def create_ppt(clean_text, output_path="output_cleaned.pptx"):
    prs = Presentation()

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Cleaned OCR Document"
    slide.placeholders[1].text = "Generated Automatically"

    # ----------------------------
    # Utility: Add a content slide
    # ----------------------------
    def add_content_slide(title, lines):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title

        body = slide.shapes.placeholders[1].text_frame
        body.clear()

        for line in lines:
            if not line.strip():
                p = body.add_paragraph()
                p.text = ""
                continue

            p = body.add_paragraph()
            p.text = line.strip()

            # Bullet formatting
            if line.startswith(("•", "-", "*")):
                p.level = 1
                p.text = line.lstrip("•-* ").strip()
            else:
                p.level = 0

            p.font.size = Pt(20)

    # ----------------------------
    # Slide splitting
    # ----------------------------
    lines = clean_text.split("\n")

    current_page = []
    max_lines_per_slide = 12  # avoid clutter

    for line in lines:
        current_page.append(line)

        if len(current_page) >= max_lines_per_slide:
            add_content_slide("Content", current_page)
            current_page = []

    if current_page:
        add_content_slide("Content", current_page)

    prs.save(output_path)
    return output_path
    print(f"\n✅ Cleaned text saved to PPTX → {output_path}")
    print(f"\n✅ Cleaned text saved to PPTX → {output_path}")   

