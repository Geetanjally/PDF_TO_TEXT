import os
import google.generativeai as genai
from pptx import Presentation
from docx import Document
from fpdf import FPDF

# --------------------------------
# CONFIG
# --------------------------------
API_KEY = "AIzaSyAzHf66I6a1uHUbC1-PnFCK6KyBUZTOJYI"
MODEL_NAME = "models/gemini-2.5-flash"

genai.configure(api_key=API_KEY)
model = genai.GenerativeModel(MODEL_NAME)

# --------------------------------
# Function: Read text file
# --------------------------------
def read_text_file(path):
    with open(path, "r", encoding="utf-8") as f:
        return f.read()

# --------------------------------
# Function: Clean/Edit with Gemini
# --------------------------------
def process_text_with_gemini(raw_text, user_instruction):
    prompt = f"""
You are a professional text editor. 
Here is the OCR extracted text:

### OCR TEXT:
{raw_text}

### USER INSTRUCTION:
{user_instruction}

### TASK:
- Modify the OCR text according to the user instruction
- Fix grammar, spellings only if needed
- Preserve structure where appropriate
- Return ONLY the edited text
"""

    response = model.generate_content(prompt)
    return response.text.strip()

# --------------------------------
# PDF / DOCX / PPT Export Functions
# --------------------------------
def export_pdf(text, output_path):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.set_font("Arial", size=12)

    for line in text.split("\n"):
        pdf.multi_cell(0, 10, line)

    pdf.output(output_path)

def export_docx(text, output_path):
    doc = Document()
    for line in text.split("\n"):
        doc.add_paragraph(line)
    doc.save(output_path)

def export_ppt(text, output_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()

    for line in text.split("\n"):
        p = tf.add_paragraph()
        p.text = line

    prs.save(output_path)

def export_all_outputs(text, folder, base_name):
    pdf_path = os.path.join(folder, base_name + ".pdf")
    docx_path = os.path.join(folder, base_name + ".docx")
    ppt_path = os.path.join(folder, base_name + ".pptx")

    export_pdf(text, pdf_path)
    export_docx(text, docx_path)
    export_ppt(text, ppt_path)

    return pdf_path, docx_path, ppt_path

# --------------------------------
# MAIN INTERACTIVE SYSTEM
# --------------------------------
if __name__ == "__main__":

    print("\n🔹 Interactive Gemini Processing System Running...\n")

    # 1️⃣ Choose input text file
    txt_path = r"G:\Project\PDF_TO_TEXT\2_OpenCV_OCR\test_output\combined_output.txt"

    raw_text = read_text_file(txt_path)
    print("\n📄 Extracted text loaded successfully!\n")

    # 2️⃣ User instruction for editing
    instruction = input("💬 Enter the instruction for Gemini: ")

    # 3️⃣ Process with Gemini
    print("\n🤖 Processing with Gemini...\n")
    final_text = process_text_with_gemini(raw_text, instruction)

    print("🧼 Final Processed Text:\n")
    print(final_text)

    # 4️⃣ Ask user if they want exports
    generate = input("\n📦 Generate PDF/DOCX/PPT? (yes/no): ").strip().lower()

    if generate == "yes":
        output_folder = r"G:\Project\PDF_TO_TEXT\5_Final_Outputs"
        os.makedirs(output_folder, exist_ok=True)

        export_all_outputs(final_text, output_folder, "Final_Output")

        print("\n🎉 All files generated successfully!")
        print(f"📂 Saved to: {output_folder}\n")

    print("✅ Done!")
