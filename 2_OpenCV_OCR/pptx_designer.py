import json
import io
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import ChartData

# --- THEME CONFIGURATION ---
# Defines the look and feel for the 3 default options
THEMES = {
    "Professional": {
        "title_font": "Arial",
        "title_color": (0, 51, 102),  # Navy Blue
        "body_font": "Calibri",
        "body_color": (89, 89, 89),   # Dark Gray
        "accent_color": (0, 51, 102),
        "design_element": True,
        "design_color": (0, 51, 102), # Navy Bar
        "design_shape": "RECTANGLE",  # Bottom Bar (The "bold" design element)
        "title_size": 32,
        "body_size": 18
    },
    "Creative": {
        "title_font": "Georgia",
        "title_color": (230, 81, 0),  # Deep Orange
        "body_font": "Gill Sans MT",
        "body_color": (40, 40, 40),   # Almost Black
        "accent_color": (75, 0, 130), # Indigo
        "design_element": True,
        "design_color": (230, 81, 0), # Orange Accent
        "design_shape": "SIDE_BAR",   # Side accent
        "title_size": 36,
        "body_size": 20
    },
    "Basic": {
        "title_font": "Calibri Light",
        "title_color": (0, 0, 0),     # Black
        "body_font": "Calibri",
        "body_color": (60, 60, 60),   # Gray
        "accent_color": (0, 0, 0),
        "design_element": False,      # Clean, no shapes
        "design_color": (255, 255, 255),
        "design_shape": None,
        "title_size": 30,
        "body_size": 18
    }
}

# --- Constants ---
# Placeholder IDs based on standard MSO types for compatibility
# The CONTENT/BODY placeholder is typically assigned the idx=2
BODY_PLACEHOLDER_IDX = 2 
DEFAULT_THEME = "Professional" # Fallback if theme_name is invalid

def _apply_theme_style(run, level, theme_cfg):
    """Applies font family, size, and color based on the selected theme."""
    font = run.font
    font.name = theme_cfg["body_font"]
    
    # Scale font size slightly smaller for deeper bullet levels
    base_size = theme_cfg["body_size"]
    font.size = Pt(base_size - (level * 2))
    
    # Apply Color
    r, g, b = theme_cfg["body_color"]
    font.color.rgb = RGBColor(r, g, b)

def _add_design_element(slide, theme_cfg):
    """Adds the visual 'bold' design element (shape) to the slide."""
    if not theme_cfg.get("design_element", False):
        return

    r, g, b = theme_cfg["design_color"]
    color = RGBColor(r, g, b)
    
    # Dimensions for a standard 16:9 slide (10in x 7.5in)
    
    if theme_cfg["design_shape"] == "RECTANGLE":
        # Professional: A sleek bar at the bottom
        left = Inches(0)
        top = Inches(7.0) 
        width = Inches(10)
        height = Inches(0.5)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background() # No outline

    elif theme_cfg["design_shape"] == "SIDE_BAR":
        # Creative: A vertical accent on the left
        left = Inches(0)
        top = Inches(1.5)
        width = Inches(0.2)
        height = Inches(5.5)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

def _add_placeholder_chart_slide(prs, title, chart_data_string, theme_cfg):
    """Creates a chart slide with themed title and a Clustered Column Chart."""
    # Blank layout is safest for manual chart insertion
    layout = prs.slide_layouts[6] 
    slide = prs.slides.add_slide(layout)
    
    # Apply Theme Design Element
    _add_design_element(slide, theme_cfg)

    # Add Title Box Manually (since we used blank layout 6)
    left = top = Inches(0.5)
    width = Inches(9.0)
    height = Inches(1.0)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    title_placeholder = txBox.text_frame
    title_placeholder.clear()
    
    p = title_placeholder.paragraphs[0]
    p.text = title
    
    # Apply Theme Font to Title
    p.font.name = theme_cfg["title_font"]
    p.font.size = Pt(theme_cfg["title_size"])
    r, g, b = theme_cfg["title_color"]
    p.font.color.rgb = RGBColor(r, g, b)


    # Chart Data Parsing and Generation
    try:
        # Expected format: [CHART:<Chart Type Title>, <Series Name>, <Category 1>:<Value 1>, ...]
        parts = chart_data_string.split(',')
        if len(parts) < 3:
            raise ValueError("Chart string must contain a series name and at least one data point.")
            
        chart_title = parts[1].strip()
        
        categories = []
        series_data = []
        
        for part in parts[2:]:
            if ':' in part:
                label, value = part.split(':', 1)
                categories.append(label.strip())
                try:
                    # Clean value before conversion (removes currency symbols, commas, etc.)
                    series_data.append(float(re.sub(r'[^\d.]', '', value)))
                except ValueError:
                    series_data.append(0.0)

        data = ChartData()
        data.categories = categories
        data.add_series(chart_title, series_data)

        # Use Clustered Column Chart
        chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
            
        # Chart position and size
        x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(5)
        
        chart = slide.shapes.add_chart(chart_type, x, y, cx, cy, data).chart
        
        # Set the chart title to the one provided in the data
        chart.has_title = True
        chart.chart_title.text_frame.text = chart_title
        
        return True
    except Exception as e:
        # Return False to indicate chart generation failure
        return False

def create_pptx_with_style(slides_data, theme_name=DEFAULT_THEME, template_data=None):
    """
    Generates a PPTX presentation from a structured JSON blueprint, applying custom 
    theming, design elements, and supporting chart generation.

    If `template_data` (bytes) is provided, it uses the custom template and disables 
    all default theme styling (fonts, colors, design shapes).
    
    Parameters:
    - slides_data (str): JSON string of the presentation blueprint (list of slides).
    - theme_name (str): The name of the default theme to apply ("Professional", "Creative", "Basic").
    - template_data (bytes, optional): Raw bytes of an uploaded custom PPTX template.
    
    Returns:
    - tuple: (BytesIO stream of PPTX file, error_message or None)
    """
    try:
        data = json.loads(slides_data)
    except:
        return None, "Failed to parse JSON blueprint."

    # --- 1. Load Presentation Object and Theme Configuration ---
    if template_data:
        # Use the custom template provided by the user
        template_stream = io.BytesIO(template_data)
        prs = Presentation(template_stream)
        # Custom templates ignore the default style settings. Set config for no custom design.
        theme_cfg = {"title_font": None, "body_font": None, "design_element": False} 
    else:
        # Use the default blank presentation and apply theme styling
        prs = Presentation() 
        theme_cfg = THEMES.get(theme_name, THEMES[DEFAULT_THEME])

    # --- 2. Iterate and Build Slides ---
    for i, slide_data in enumerate(data):
        slide_title = slide_data.get('title', f"Slide {i+1}")
        slide_content = slide_data.get('content', [])
        
        # Check for Chart Slide/Content
        chart_placeholder = next((p.strip() for p in slide_content if p.strip().upper().startswith('[CHART:')), None)
        remaining_content = [p for p in slide_content if not p.strip().upper().startswith('[CHART:')]

        # --- A. Generate Chart Slide ---
        if chart_placeholder:
            # We call the helper function to draw the chart and themed title manually
            # The data is parsed from the string: [CHART:Chart Title, Series Name, Categ1:Value1, ...]
            # Remove [CHART: and ]
            _add_placeholder_chart_slide(prs, slide_title, chart_placeholder[7:-1], theme_cfg)
            continue 

        # --- B. Generate Standard Content Slide ---
        
        # Use layout 1 (Title and Content) or layout 5 (Title only) as a reliable base
        layout_index = 1 if remaining_content else 5 # Title and Content (1) or Title Only (5)
        
        try:
            layout = prs.slide_layouts[layout_index] 
        except IndexError:
            layout = prs.slide_layouts[0] # Fallback to Title Slide (Index 0 is always safe)

        slide = prs.slides.add_slide(layout)
        
        # Apply Design Element ONLY if no template was provided
        if not template_data:
            _add_design_element(slide, theme_cfg)
        
        # 1. Title Styling
        try:
            slide.shapes.title.text = slide_title
            # Apply Theme Title Font/Color ONLY if no template was provided
            if not template_data:
                tp = slide.shapes.title.text_frame.paragraphs[0]
                tp.font.name = theme_cfg["title_font"]
                tp.font.size = Pt(theme_cfg["title_size"])
                tr, tg, tb = theme_cfg["title_color"]
                tp.font.color.rgb = RGBColor(tr, tg, tb)
        except:
            # Handle slides without a title placeholder gracefully
            pass
            
        # 2. Body Content (Only if layout index 1 was used and there is content)
        if layout_index == 1 and remaining_content:
            try:
                # --- FIX: Robustly find the content placeholder by its MSO idx (2) ---
                body_placeholder = None
                for shape in slide.shapes:
                    # Check if it's a placeholder and its MSO index is 2 (Body Content)
                    if shape.is_placeholder and shape.placeholder_format.idx == BODY_PLACEHOLDER_IDX:
                        body_placeholder = shape
                        break

                if not body_placeholder:
                    # Fallback to using the second item in the list if the robust search failed
                    if len(slide.shapes.placeholders) > 1:
                        # Attempt the most common index for Body content (Index 1)
                        body_placeholder = slide.shapes.placeholders[1]
                    else:
                        raise ValueError("No suitable body placeholder found.")

                # Proceed with content insertion
                body = body_placeholder.text_frame
                body.clear()
                
                for point in remaining_content:
                    processed_point = point.strip()
                    level = 0
                    
                    # Markdown-like level parsing
                    if processed_point.startswith('***'):
                        level = 2
                        processed_point = processed_point.lstrip('***').strip()
                    elif processed_point.startswith('**'):
                        level = 1
                        processed_point = processed_point.lstrip('**').strip()
                    elif processed_point.startswith('*'):
                        level = 0
                        processed_point = processed_point.lstrip('*').strip()

                    p = body.add_paragraph()
                    run = p.add_run() 
                    run.text = processed_point
                    
                    # Apply Theme Body Font/Color/Size ONLY if no template was provided
                    if not template_data:
                        _apply_theme_style(run, level, theme_cfg)
                    
                    # Paragraph spacing and bullet level
                    p.space_after = Pt(10)
                    p.level = level

            except Exception as e:
                print(f"Warning: Could not populate content for slide {i+1} due to text frame issue. Error: {e}")
                # Fallback to a plain textbox if placeholder insertion fails
                if remaining_content:
                    # Fallback: If no standard content placeholder is found, dump the content into a new textbox
                    left = top = Inches(1)
                    width = Inches(8.5)
                    height = Inches(5.5)
                    
                    # Position the fallback textbox away from the assumed title position
                    txBox = slide.shapes.add_textbox(left, top + Inches(1.5), width, height) 
                    tf = txBox.text_frame
                    tf.clear()
                    tf.vertical_anchor = MSO_ANCHOR.TOP
                    
                    for point in remaining_content:
                        # Replicate bullet structure in the fallback
                        processed_point = point.strip()
                        level = 0
                        if processed_point.startswith('***'):
                            level = 2
                            processed_point = processed_point.lstrip('***').strip()
                        elif processed_point.startswith('**'):
                            level = 1
                            processed_point = processed_point.lstrip('**').strip()
                        elif processed_point.startswith('*'):
                            level = 0
                            processed_point = processed_point.lstrip('*').strip()
                        
                        p = tf.add_paragraph()
                        run = p.add_run()
                        run.text = processed_point
                        p.level = level
                        
                        if not template_data:
                            _apply_theme_style(run, level, theme_cfg)


    # Save
    pptx_stream = io.BytesIO()
    prs.save(pptx_stream)
    pptx_stream.seek(0)
    return pptx_stream, None