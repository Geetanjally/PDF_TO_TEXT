import os
import io
import re
import time
import json
import base64
import numpy as np
import cv2
import fitz  # PyMuPDF
from pptx.dml.color import RGBColor  # REQUIRED for coloring image prompts

from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE

# Import the necessary Google GenAI libraries
from google.genai import Client
from google.genai.types import Image, GenerateContentConfig, Part
from google.api_core import exceptions

# --- Configuration Constants ---
client = Client(api_key=os.getenv("GEMINI_API_KEY"))

MODEL_NAME = "models/gemini-flash-latest"
OCR_MODEL_NAME = "models/gemini-flash-latest"

# ----------------------------------------------------------------------
# 1. PDF/Image Extraction (Modified for in-memory processing for Streamlit)
# ----------------------------------------------------------------------

def extract_text_from_pdf(pdf_bytes, dpi=200):
    """
    Extracts text from PDF pages, and renders pages as images for OCR if text is sparse.
    Returns a list of tuples: [(page_number, extracted_text, base64_image)]
    
    Ensures iteration over ALL pages and saves images (PNG) for pages with sparse digital text.
    """
    
    # Check for PyMuPDF
    if not hasattr(fitz, 'open'):
        raise ImportError("PyMuPDF (fitz) is not installed. Please install with 'pip install pymupdf'")

    # Convert bytes to a PyMuPDF document
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    results = []
    num_pages = len(doc)
    
    print(f"✅ Opened PDF | Total pages: {num_pages}")

    try:
        # Iterate over all pages detected by PyMuPDF to ensure all pages are processed.
        for i, page in enumerate(doc):
            page_num = i + 1
            
            # 1. Attempt to extract text directly from the PDF
            raw_pdf_text = page.get_text()
            
            # 2. Render the page to an image (in memory)
            # Use matrix=fitz.Matrix(dpi/72, dpi/72) for consistent DPI across platforms
            matrix = fitz.Matrix(dpi/72, dpi/72)
            pix = page.get_pixmap(matrix=matrix, alpha=False)
            
            # Convert pixmap to image bytes (PNG is lossless and safe for OCR)
            # We must use 'png' mimeType later if we use this format.
            img_data = pix.tobytes(output="png") 
            base64_img = base64.b64encode(img_data).decode('utf-8')
            
            # 3. Decide which data structure to save
            
            # Threshold Check: If digital text has a decent amount of content, use it directly.
            # Strips whitespace and normalizes it to count meaningful characters.
            meaningful_text_length = len(re.sub(r'\s+', '', raw_pdf_text))

            # Threshold: If meaningful text > 250 characters, assume digital text is sufficient.
            if meaningful_text_length > 250:
                print(f"📄 Page {page_num}: Digital text found, skipping image OCR. Content Length: {meaningful_text_length}")
                # Use raw_pdf_text and set base64_img to None (to avoid unnecessary OCR)
                results.append((page_num, raw_pdf_text.strip(), None)) 
            else:
                print(f"🖼️ Page {page_num}: Sparse text found. Image saved for OCR. Content Length: {meaningful_text_length}")
                # Use None for text and keep base64_img for OCR
                results.append((page_num, None, base64_img)) 

    finally:
        # Crucial: Close the document to release resources
        doc.close()
        
    return results

# ----------------------------------------------------------------------
# 2. OCR Classification
# ----------------------------------------------------------------------

def is_image_digital(img_bytes):
    """
    Classifies an image as digital (printed) or handwritten using basic CV checks.
    For Streamlit, we work with image bytes/base64, not file paths.
    """
    # Guard against missing OpenCV dependency
    if cv2 is None or np is None:
        print("[⚠️] OpenCV/Numpy not available for classification. Assuming handwritten/fuzzy.")
        return False

    try:
        # Convert bytes to a numpy array for OpenCV
        nparr = np.frombuffer(img_bytes, np.uint8)
        img = cv2.imdecode(nparr, cv2.IMREAD_COLOR)

        if img is None:
            print("[⚠️] Could not decode image bytes.")
            return False

        gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)

        # 1️⃣ Quick check: edge density 
        edges = cv2.Canny(gray, 50, 150)
        edge_density = np.sum(edges > 0) / edges.size

        # Simple thresholding: high edge density (e.g., > 0.01) suggests sharp, printed text
        is_digital = edge_density > 0.01
        
        print(f"[ℹ️] Edge Density: {edge_density:.4f} -> {'Digital' if is_digital else 'Handwritten/Fuzzy'}")
        
        return is_digital

    except Exception as e:
        print(f"❌ Classification Error: {e}")
        # Default to Gemini for safety if classification fails
        return False 

# ----------------------------------------------------------------------
# 3. OCR Engines (Tesseract & Gemini)
# ----------------------------------------------------------------------

def extract_text_tesseract(img_bytes, timeout=10, max_retries=2):
    """
    Tesseract OCR (for printed/digital text) - Note: requires local tesseract installation.
    Since Tesseract is usually not available in cloud environments, this acts as a placeholder
    and returns an empty string, forcing the flow to rely on Gemini Vision.
    """
    print("⚠️ Tesseract (local OCR) is generally not supported in cloud environments like Streamlit. Bypassing and relying on Gemini.")
    return "" 

def extract_text_gemini(data, api_key, is_base64=True, max_retries=5):
    """
    Gemini OCR with built-in Retry Logic for 429 errors.
    """
    prompt = ("Extract ALL text accurately. Preserve formatting, steps, bullet points, equations, "
        "indentation, tables, and line breaks. Do NOT summarize. Return ONLY the raw text.")
    img_data = base64.b64decode(data) if is_base64 else data
    img_part = Part.from_bytes(data=img_data, mime_type='image/png')

    config = GenerateContentConfig(temperature=0, max_output_tokens=8192)

    delay = 5  # Initial wait time in seconds
    for attempt in range(max_retries):
        try:
            print(f"✨ Gemini OCR attempt {attempt + 1}/{max_retries}")
            response = client.models.generate_content(
                model=OCR_MODEL_NAME,
                config=config,
                contents=[img_part, prompt]
            )
            return response.text.strip()

        except exceptions.ResourceExhausted as e:
            # This handles the 429 error specifically
            if attempt < max_retries - 1:
                print(f"⚠️ Rate limit hit. Waiting {delay} seconds before retry...")
                time.sleep(delay)
                delay *= 2  # Wait longer each time (5s, 10s, 20s...)
            else:
                return f"[ERROR] Gemini OCR failed after {max_retries} attempts due to rate limiting."
        except Exception as e:
            print(f"❌ Other error: {e}")
            break
            
    return ""

# ----------------------------------------------------------------------
# 4. Text Cleaning and Structuring (CHUNKED VERSION)
# ----------------------------------------------------------------------

def _clean_raw_text(text):
    """
    Helper function to clean raw OCR text.
    Removes excessive whitespace and normalizes line breaks.
    """
    # Remove excessive blank lines (more than 2 newlines in a row)
    text = re.sub(r'\n\s*\n\s*\n+', '\n\n', text)
    # Remove trailing/leading whitespace from each line
    text = '\n'.join(line.strip() for line in text.split('\n'))
    return text.strip()

def chunk_pages(pages_list, chunk_size=5):
    """Utility to split page list into smaller manageable groups."""
    for i in range(0, len(pages_list), chunk_size):
        yield pages_list[i:i + chunk_size]

def clean_chunk_with_gemini(raw_text, api_key, part_no, total_parts):
    """
    Uses Gemini to clean a specific chunk with awareness of its position.
    """
    if not api_key:
        return "[ERROR] API Key is missing."

    # Context-aware prompt to maintain continuity
    prompt = f"""
You are a professional text cleaner. This is PART {part_no} of {total_parts} of a larger study document.

STRICT INSTRUCTIONS:
- If this is NOT Part 1, skip the general introduction.
- If this is NOT the last Part, skip the concluding summary.
- Fix OCR errors and preserve ALL technical details and bullet points.
- Maintain Markdown structure (# for topics, ## for sections).
- Strategically insert [Image of X] tags for complex concepts that would benefit from visual aids (diagrams, charts, illustrations).

OCR Text from Part {part_no}:
{raw_text}

Return ONLY the cleaned Markdown.
"""
    
    gen_config = GenerateContentConfig(temperature=0.2, max_output_tokens=8192)

    try:
        response = client.models.generate_content(
            model=MODEL_NAME,
            contents=[prompt],
            config=gen_config
        )
        return response.text.strip()
    except Exception as e:
        print(f"❌ Chunk {part_no} Cleaning Error: {e}")
        return f"\n[Error cleaning Part {part_no}]\n"

# ----------------------------------------------------------------------
# 5. Core Pipeline Function (Scalable & Robust)
# ----------------------------------------------------------------------

def process_document_to_cleaned_text(pdf_file_bytes, api_key):
    """
    Handles PDF -> Page Extraction -> Chunking -> Parallelized-style Cleaning.
    """
    # 1. Extract raw content from every page
    page_results = extract_text_from_pdf(pdf_file_bytes)
    all_raw_pages = []
    
    print(f"Starting extraction for {len(page_results)} pages...")

    for page_num, digital_text, base64_img in page_results:
        content = ""
        if digital_text:
            content = digital_text.strip()
        elif base64_img:
            # Fallback to OCR
            ocr_result = extract_text_gemini(base64_img, api_key)
            content = ocr_result.strip()
            # Pause briefly to respect rate limits for OCR calls
            time.sleep(1) 
        
        if content:
            all_raw_pages.append(content)

    if not all_raw_pages:
        return None, "[ERROR] No text could be extracted from this PDF."

    # 2. Split pages into chunks (e.g., 5 pages per Gemini call)
    PAGES_PER_CHUNK = 5
    page_chunks = list(chunk_pages(all_raw_pages, PAGES_PER_CHUNK))
    total_parts = len(page_chunks)
    cleaned_chunks = []

    print(f"Processing {total_parts} chunks for cleaning...")

    # 3. Clean each chunk sequentially
    for idx, chunk in enumerate(page_chunks, start=1):
        # Combine pages in this chunk
        chunk_raw_text = "\n\n---\n\n".join(chunk)
        chunk_raw_text = _clean_raw_text(chunk_raw_text)

        print(f"🧹 Cleaning chunk {idx}/{total_parts}...")
        cleaned_text = clean_chunk_with_gemini(chunk_raw_text, api_key, idx, total_parts)
        cleaned_chunks.append(cleaned_text)

        # CRITICAL: Sleep between chunks to avoid 429 Rate Limit errors
        time.sleep(4)

    # 4. Merge all cleaned parts into the final document
    final_cleaned_text = "\n\n".join(cleaned_chunks)
    
    return final_cleaned_text, None

# ----------------------------------------------------------------------
# 6. Final Output Generation (Memory-based for Streamlit)
# ----------------------------------------------------------------------

def create_pptx_from_markdown(markdown_text):
    """
    Generates a PPTX file in memory from structured Markdown text.
    Converts [Image of X] tags into colored prompts for manual image insertion.
    """
    prs = Presentation() 
    # Widescreen 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Split slides by Markdown headings (# or ##)
    # Added \s to ensure it only splits on actual headers
    slides_data = re.split(r'(?=\n#+\s)', markdown_text)
    slides_data = [s.strip() for s in slides_data if s.strip()]

    title_slide_layout = prs.slide_layouts[0] 
    content_slide_layout = prs.slide_layouts[1]

    for i, slide_text in enumerate(slides_data):
        lines = slide_text.split('\n')
        title_line = lines[0].strip()
        
        # Clean the title (remove # symbols)
        title_match = re.match(r'(#+)\s*(.*)', title_line)
        title = title_match.group(2).strip() if title_match else title_line
        
        # Add slide
        slide = prs.slides.add_slide(title_slide_layout if i == 0 else content_slide_layout)
        slide.shapes.title.text = title
        
        # Get the body text frame
        try:
            body_shape = slide.shapes.placeholders[1]
            body = body_shape.text_frame
            body.clear()  # Clear any default placeholder text
        except:
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
            body = txBox.text_frame
        
        body.word_wrap = True
        content_lines = lines[1:]

        for line in content_lines:
            line = line.strip()
            if not line:
                continue

            p = body.add_paragraph()
            
            # --- 1. HANDLE [Image of X] TAGS ---
            # This looks for: [Image of X] where X is any text
            img_match = re.search(r'\[Image of (.+?)\]', line, re.IGNORECASE)
            
            if img_match:
                image_topic = img_match.group(1).strip()
                p.text = f"🖼️ [PROMPT: {image_topic}]"
                p.font.bold = True
                p.font.color.rgb = RGBColor(0, 102, 204)  # Blue color for visibility
                continue

            # --- 2. HANDLE BULLETS AND FORMATTING ---
            if re.match(r'^[\*•-]\s', line):
                p.text = re.sub(r'^[\*•-]\s', '', line).strip()
                p.level = 0
            elif line.startswith('  ') or line.startswith('\t') or re.match(r'^[\*•-]{2,}', line):
                p.text = line.lstrip('*•- \t').strip()
                p.level = 1
            elif line.startswith('###'):
                p.text = line.lstrip('#').strip()
                p.font.bold = True
                p.level = 0
            else:
                p.text = line
                p.level = 0

            p.font.size = Pt(18) 

    prs_io = io.BytesIO()
    prs.save(prs_io)
    prs_io.seek(0)
    return prs_io

def create_docx_from_markdown(markdown_text):
    """
    Generates a DOCX file in memory from structured Markdown text.
    """
    doc = Document()
    
    for line in markdown_text.split('\n'):
        line = line.strip()
        if not line:
            continue
            
        if line.startswith('###'):
            doc.add_heading(line.lstrip('#').strip(), level=3)
        elif line.startswith('##'):
            doc.add_heading(line.lstrip('#').strip(), level=2)
        elif line.startswith('#'):
            doc.add_heading(line.lstrip('#').strip(), level=1)
        elif line.startswith(('*', '-', '•')):
            doc.add_paragraph(line.lstrip('*-• ').strip(), style='List Bullet')
        else:
            doc.add_paragraph(line)
            
    # Save the document to a byte stream
    doc_io = io.BytesIO()
    doc.save(doc_io)
    doc_io.seek(0)
    return doc_io
    
def create_markdown_report(markdown_text):
    """
    Returns the final structured Markdown text as a file-like object.
    """
    # Ensure text is encoded to bytes for a file download
    markdown_io = io.BytesIO(markdown_text.encode('utf-8'))
    markdown_io.seek(0)
    return markdown_io