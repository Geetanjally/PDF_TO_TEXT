import os
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.platypus import SimpleDocTemplate, Paragraph
from reportlab.lib.styles import getSampleStyleSheet

# -------------------------------------------------------------
# 1️⃣  Generate PDF
# -------------------------------------------------------------
def generate_pdf(text, output_path):
    styles = getSampleStyleSheet()
    style = styles["Normal"]
    pdf = SimpleDocTemplate(output_path)

    flow = [Paragraph(text.replace("\n", "<br/>"), style)]
    pdf.build(flow)

    print("📄 PDF created:", output_path)


# -------------------------------------------------------------
# 2️⃣  Generate Word Document
# -------------------------------------------------------------
def generate_docx(text, output_path):
    doc = Document()
    for line in text.split("\n"):
        doc.add_paragraph(line)
    doc.save(output_path)
    print("📝 DOCX created:", output_path)


# -------------------------------------------------------------
# 3️⃣  Generate PPT (slides based on headings & bullets)
# -------------------------------------------------------------
def generate_ppt(text, output_path):
    prs = Presentation()

    # Split into slides using double newlines
    slides = text.split("\n\n")

    for slide_text in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content

        lines = slide_text.strip().split("\n")
        
        title = lines[0] if len(lines) > 0 else "Slide"
        content = "\n".join(lines[1:]) if len(lines) > 1 else ""

        slide_title = slide.shapes.title
        slide_body = slide.placeholders[1]

        slide_title.text = title
        slide_body.text = content

    prs.save(output_path)
    print("📊 PPT created:", output_path)


# -------------------------------------------------------------
# 4️⃣  Master function
# -------------------------------------------------------------
def export_all_outputs(text, output_folder, base_name):
    os.makedirs(output_folder, exist_ok=True)

    pdf_path = os.path.join(output_folder, f"{base_name}.pdf")
    docx_path = os.path.join(output_folder, f"{base_name}.docx")
    ppt_path = os.path.join(output_folder, f"{base_name}.pptx")

    generate_pdf(text, pdf_path)
    generate_docx(text, docx_path)
    generate_ppt(text, ppt_path)

    print("\n✅ All files generated successfully!\n")
